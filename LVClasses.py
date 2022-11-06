import math
from math import sqrt, pow, pi, cos, sin, tan, acos, asin, atan, radians, degrees, exp, log
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches

class Mission:
    """Declares the Mission parameters inside a Mission object. Uses function getTrajReqs to find the delta-v trajectory requirements needed to size the LV"""
    pass
    # CONSTANTS OF EARTH
    g = 0.00980065  # gravity (km/s^2)
    mu_E = 398600  # gravitational parameter of earth (km^3/s)
    r_E = 6378  # radius of earth (km)
    v_equator = 0.4651 # equatorial velocity in (km/s)
     
    # Delta-V Trajectory Requirements Variable
    dV_reqs = '\'dv_reqs\' is not yet initialized, call the set method \'set_dV_reqs\''

    # CONSTRUCTOR
    def __init__(self, Mission_type, recovery, losses_gravity, drag_loss, launch_site):
        #Mission_type Construct an instance of this class
        #   Detailed explanation goes here
        self.input = [Mission_type, recovery, losses_gravity, drag_loss, launch_site]

    # Delta-V Design Function
    def set_dV_reqs(self):
        #print(self.input)
        if self.input[0] == 'One':
            delta_plane = 0
            inc = radians(60)
            h_a = 500 # apoapsis altitude (km)
        elif self.input[0] =='Two':
            delta_plane = radians(10) # plane change of 10 degrees (rad)
            inc = radians(98) # inclination (radians)
            h_a = 550 # apoapsis altitude (km)

        h_p = 200 # periapsis altitude (km)
        if self.input[4] == 'Kodiak':
            lat = radians(57.79) # latitude (radians)
   
        ### Orbital Calculations
        #  Orbital velocities
        r_p = self.r_E + h_p # periapsis radius (km)
        r_a = self.r_E + h_a # apoapsis radius (km)
        a = (r_p + r_a)/2 # semi-major axis (km)
        #print(r_p, r_a, a)

        v_p = sqrt(2*self.mu_E*(1/r_p - 1/(2*a))) # periapsis velocity (km/s)
        v_a = sqrt(2*self.mu_E*(1/r_a - 1/(2*a))) # apoapsis velocity (km/s)
        v_c = sqrt(self.mu_E/r_a) # circular velocity at h_a (km/s)
        v_LS = self.v_equator * cos(lat) # launch site velocity (km/s)
        #print(v_p, v_a, v_c, v_LS)

    

        # Angles
        if inc > pi/2:
            aux = pi - inc # launch window auxiliary angle (rad)
        elif inc < pi/2:
            aux = inc # launch window auxiliary angle (rad)
        flt_path = asin(cos(aux)/cos(lat)) # flight path angle (rad)
        #print(aux, flt_path)

        # Azimuth Angle and Burnout Velocities
        if inc > pi/2:
            azimuth = pi + flt_path # azimuth angle (rad)
            v_BO_S = -v_p*cos(flt_path)*cos(azimuth) # South burnout velocity (km/s)
            v_BO_E = -v_p*cos(flt_path)*sin(azimuth) # East burnout velocity (km/s)
            v_BO_Z = v_p*sin(flt_path) # Zenith burnouth velocity (km/s)

        elif inc < pi/2:
            azimuth = flt_path
            v_BO_S = -v_p*cos(flt_path)*cos(azimuth) # South burnout velocity (km/s)
            v_BO_E = v_p*cos(flt_path)*sin(azimuth) # East burnout velocity (km/s)
            v_BO_Z = v_p*sin(flt_path) # Zenith burnouth velocity (km/s)

        #print(v_BO_S, v_BO_E, v_BO_Z)
        v_N_S = v_BO_S # South needed delta-v (km/s)
        v_N_E = v_BO_E - v_LS # East needed delta-v (km/s)
        v_N_Z = v_BO_Z # Zenith needed delta-v (km/2)
        #print(v_N_S, v_N_E, v_N_Z)

        dv_maneuvers = 0
        if self.input[1]:
            dv_maneuvers += 0.343 # landing dv burn (km/s)
        else: dv_maneuvers +=0

        dv_N = sqrt(pow(v_N_S,2) + pow(v_N_E,2) + pow(v_N_Z,2)) # total delta-v needed (km/s)
        grav_loss = 0
        if (type(self.input[2]) == str) & (self.input[2] == '80% gravity loss'):
            grav_loss = 0.8*sqrt(2*self.mu_E*h_p/((h_p+self.r_E)*self.r_E)) # 80# gravity loss eqn (km/s)
        elif (type(self.input[2]) == int) | (type(self.input[2]) == float):
            grav_loss = self.input[1]
        apo_kick = v_c - v_a # apoapsis kick burn (km/s)

        if self.input[0] == 'One':
            dv_design = dv_N + grav_loss + self.input[3] + apo_kick + dv_maneuvers # delta-v design is the total delta v required (km/s)
            #print(dv_N, grav_loss, drag_loss, apo_kick, dv_maneuvers, dv_design)
            self.dV_reqs = [dv_N, dv_design, 0, grav_loss, self.input[3], apo_kick, dv_maneuvers]
        elif self.input[0] =='Two':
            dv_plane = 2*v_a*sin(delta_plane/2) # delta-v needed for plane-change (km/s)
            dv_design = dv_N + grav_loss + self.input[3] + apo_kick + dv_maneuvers + dv_plane # delta-v design is the total delta v required (km/s)
            #print(dv_N, grav_loss, drag_loss, apo_kick, dv_maneuvers, dv_plane, dv_design)
            self.dV_reqs = [dv_N, dv_design, dv_plane, grav_loss, self.input[3], apo_kick, dv_maneuvers]

    def generateFamilies(LV, percent_dV, previous=0): # takes a launch vehicle object and a list of percent_dV distribution for each step except the last, i.e. [0.5], [0.5, 0.3]
        pass
        #create vector to store percent delta-v's, for each step
        percent_dV.app (1-sum(percent_dV)) # app  the percent dV from the last step 
        LV_dV_reqs = LV.Mission.dV_reqs # get dV trajectory requirements from Mission  
        stage_losses = [] # initialize list of stage_losses for each step 
        staging_speed = [] # initialize list of staging speeds for each step
        required_dv = [] # initialize list of required delta-v's for each step
        MR = [] # initialize list of mass ratios for each step
        m_p = [] # initialize list of propellant masses for each step
        m_s = [] # initialize list of structural masses for each step
        m_step = [] # initialize list of total step mass for each step
        for i in range(LV.num_steps):
            # app  values to each list by equations in textbook on p.283 & p.284
            if(i != LV.num_steps-1):
                stage_losses.app ((LV_dV_reqs[3]+LV_dV_reqs[4]+LV_dV_reqs[6])/(LV.num_steps-1)) # sums together and splits gravity, drag, and maneuver losses evenly between all steps other than the last step
            else:
                stage_losses.app (LV_dV_reqs[2] + LV_dV_reqs[5]) # sums together the dV_plane change and dV_apo_kick for stage 2 'losses'
            staging_speed.app (LV_dV_reqs[0]*percent_dV[i])
            required_dv.app (stage_losses[i] + staging_speed[i])
            MR.app (required_dv[i]/(LV.engine_Isps[i]*LV.Mission.g)) # mass ratio is MR = e^(dv_required/(Isp*g0))
            m_p.app (LV.PL*((MR[i]-1)*(1-LV.sigmas[i])/(1-LV.sigmas[i]*MR[i])))
            m_s.app (m_p[i]*LV.sigmas[i]/(1-LV.sigmas[i]))
            if (i != LV.num_step-1):
                m_step.app (m_p[i] + m_s[i])
            else: 
                m_step.app (LV.PL*MR[i]*(1-LV.sigmas[i])/(1-MR[i]*LV.sigmas[i]))
        m_00 = sum(m_step) # Leave blank for now
    def print(self):
        print('Here are the mission requirements for Mission ' + self.input[0] + ', where the recovery is ' + str(self.input[1]) + ' and the launch site is ' + str(self.input[4]))
        print("The delta-v needed is " + str(self.dV_reqs[0]) + " km/s")
        print("The delta-v design is " + str(self.dV_reqs[1]) + " km/s")
        print("The delta-v plane change is " + str(self.dV_reqs[2]) + " km/s")
        print("The delta-v gravity loss is " + str(self.dV_reqs[3]) + " km/s")
        print("The delta-v drag loss is " + str(self.dV_reqs[4]) + " km/s")
        print("The delta-v apo-kick is " + str(self.dV_reqs[5]) + " km/s")
        print("The delta-v maneuvers (entry and landing burn) is " + str(self.dV_reqs[6]) + " km/s")

    def GoalSeek(fun,goal,x0,fTol=0.0001,MaxIter=1000):
    # Goal Seek function of Excel
    #   via use of Line Search and Bisection Methods

    # Inputs
    #   fun     : Function to be evaluated
    #   goal    : Expected result/output
    #   x0      : Initial estimate/Starting point

    # Initial check
        if fun(x0)==goal:
            print('Exact solution found')
            return x0

        # Line Search Method
        step_sizes=np.logspace(-1,4,6)
        scopes=np.logspace(1,5,5)

        vFun=np.vectorize(fun)

        for scope in scopes:
            break_nested=False
            for step_size in step_sizes:

                cApos=np.linspace(x0,x0+step_size*scope,int(scope))
                cAneg=np.linspace(x0,x0-step_size*scope,int(scope))

                cA=np.concatenate((cAneg[::-1],cApos[1:]),axis=0)

                fA=vFun(cA)-goal

                if np.any(np.diff(np.sign(fA))):

                    index_lb=np.nonzero(np.diff(np.sign(fA)))

                    if len(index_lb[0])==1:

                        index_ub=index_lb+np.array([1])

                        x_lb=np.asscalar(np.array(cA)[index_lb][0])
                        x_ub=np.asscalar(np.array(cA)[index_ub][0])
                        break_nested=True
                        break
                    else: # Two or more roots possible

                        index_ub=index_lb+np.array([1])

                        print('Other solution possible at around, x0 = ', np.array(cA)[index_lb[0][1]])

                        x_lb=np.asscalar(np.array(cA)[index_lb[0][0]])
                        x_ub=np.asscalar(np.array(cA)[index_ub[0][0]])
                        break_nested=True
                        break

            if break_nested:
                break
        if not x_lb or not x_ub:
            print('No Solution Found')
            return

        # Bisection Method
        iter_num=0
        error=10

        while iter_num<MaxIter and fTol<error:
        
            x_m=(x_lb+x_ub)/2
            f_m=fun(x_m)-goal

            error=abs(f_m)

            if (fun(x_lb)-goal)*(f_m)<0:
                x_ub=x_m
            elif (fun(x_ub)-goal)*(f_m)<0:
                x_lb=x_m
            elif f_m==0:
                print('Exact spolution found')
                return x_m
            else:
                print('Failure in Bisection Method')
        
            iter_num+=1

        return x_m


            

class LaunchVehicle(Mission) :
    pass
    def __init__(self, name, body_material, num_steps, engine_Isps, sigmas, PL, Mission):
        self.name = name
        self.body_material = body_material
        if body_material == "Aluminum 6061":
            self.rho_body = 2700  # kg/m^3 LV body material
        self.num_steps = num_steps
        self.engine_Isps = engine_Isps
        self.sigmas = sigmas
        self.PL = PL
        self.Mission = Mission



    def addMassesLV(self): # Temporary method. the m_gross and m_p parameters will be found using the given parameters in the init method by generating LVFamilies
        self.m_gross = 6414.5 # Gross lift off mass
        self.m_p = [5334.7, 474.09] # step propellant masses
        self.m_0 = [5798.6, 615.983] # step masses

    def initSteps(self, listOfSteps):
        self.listOfSteps = listOfSteps

    def initInterstages(self): # also sets the distance which the 1st step's engines protrudes out of the aft_skirt, and adds interstage to Step object
        #self.listOfInterstages = []
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            temp_interstage = []
            if i == len(self.listOfSteps)-1: # if last step: last step doesn't have interstage
                temp_interstage.append(0)
                temp_interstage.append(0)
            else: # if not last step
                if step.propulsion == 'Liquid':
                    temp_interstage.append(step.dome_f[0] + self.listOfSteps[i+1].L_n + step.r/4) # gap of r/4 between the current step's tank dome and upper step's nozzle exit
                elif step.propulsion == 'Solid':
                    temp_interstage.append(self.listOfSteps[i+1].L_n + step.r/4) # gap of r/4 between the current step's srm_casing and upper step's nozzle exit
                temp_interstage.append(temp_interstage[i]*step.circumference)
            
            #self.listOfInterstages.append(interstage)
            step.interstage = temp_interstage # set temp_interstage as the step's 'interstage' attribute

        #self.listOfSteps[0].eng_protrude_dist = self.listOfSteps[0].L_n/2
        

    def massMoments(self):
        payload_items = ['PLF', 'Payload', 'PAF']
        step_items_liquid = ['Forward Skirt', 'Avionics', 'Wiring', 'Fuel Dome Top', 'Fuel Cylinder', 'Fuel Dome Bottom', 'Fuel Insulation', 'Fuel Residual', 'Intertank', 'Ox Dome Top', 'Ox Cylinder', 'Ox Dome Bottom', 'Ox Insulation', 'Ox Residual', 'Aft Skirt', 'Thrust Structure', 'Gimballs', 'Engines', 'Fuel', 'Oxidizer']
        step_items_solid =['Forward Skirt', 'Avionics', 'Wiring', 'Solid Propellant Casing', 'Solid Propellant Residual', 'Aft Skirt', 'Thrust Structure', 'Gimballs', 'Nozzle', 'Solid Propellant']
        
        #num_items = len(payload_items) + self.num_steps * len(step_items)
        df_temp = pd.DataFrame(columns=['Item', 'Height (m)', 'Mass (kg)', 'Distance (m)', 'Moment (kg*m)', 'Thickness (m)', 'Distance from CM (m)', 'J0 (kg m^2)', 'm*CM^2 (kg m^2)', 'Jpitch/yaw', 'Jroll'],
                          index=range(len(payload_items)))
        #df = pd.DataFrame(columns=['Item', 'Height (m)', 'Mass (kg)', 'Distance (m)', 'Moment (kg*m)', 'Thickness (m)', 'Distance from CM (m)', 'J0 (kg m^2)', 'm*CM^2 (kg m^2)', 'Jpitch/yaw', 'Jroll'],
        #                  index=range(num_items))

        self.df = self.appendItems(df_temp, payload_items, step_items_liquid, step_items_solid)
        self.initHeights(self.df, len(payload_items), len(step_items_liquid), len(step_items_solid))
        self.initThicknesses(self.df, len(payload_items), len(step_items_liquid), len(step_items_solid))
        self.initMasses(self.df, len(payload_items), len(step_items_liquid), len(step_items_solid))
        self.initDistances(self.df, len(payload_items), len(step_items_liquid), len(step_items_solid))
        print(self.df)
        print(self.df[['Item', 'Thickness (m)']])
        print(self.df[['Item', 'Distance (m)']])

    def appendItems(self, df, payload_items, step_items_liquid, step_items_solid):# initialize 'Item' names
        for i in range(len(payload_items)): # initialize payload item names
            df['Item'][i] = payload_items[i]
        for i in range(len(self.listOfSteps)): # initialize step item names
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid': # if liquid step
                for j in range(len(step_items_liquid)):
                    new_row = {'Item': step_items_liquid[j] + ' ' + str(step.step_num)}
                    df = df.append(new_row, ignore_index=True)
                    
            elif step.propulsion == 'Solid': # if solid step
                for j in range(len(step_items_solid)):
                    new_row = {'Item': step_items_solid[j] + ' ' + str(step.step_num)}
                    df = df.append(new_row, ignore_index=True)
        return df

    def initHeights(self, df, num_PL, num_liquid, num_solid ):
        #print(self.listOfSteps)
        df['Height (m)'][0] = self.listOfSteps[len(self.listOfSteps)-1].fairing[0] # fairing length
        df['Height (m)'][1] = self.listOfSteps[len(self.listOfSteps)-1].fairing[0] # payload length
        df['Height (m)'][2] = " " # PAF length
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid':
                if i < len(self.listOfSteps) - 1: # if not the last step
                    df['Height (m)'][num_PL + num_liquid*i] = step.interstage[0] # Interstage 1 (still called fwd skirt)
                
                elif i == len(self.listOfSteps) - 1: #else if last step
                    df['Height (m)'][num_PL + num_liquid*i] = step.fwd_skirt[0] # fwd skirt

                df['Height (m)'][num_PL + num_liquid*i + 1] = " "                                                   # Avionics
                df['Height (m)'][num_PL + num_liquid*i + 2] = step.total_length                                     # Wiring (Total length of step)
                df['Height (m)'][num_PL + num_liquid*i + 3] = step.dome_f[0]                                        # Fuel Dome Top
                df['Height (m)'][num_PL + num_liquid*i + 4] = step.cyl_f[0]                                         # Fuel Cylinder
                df['Height (m)'][num_PL + num_liquid*i + 5] = step.dome_f[0]                                        # Fuel Dome Bottom
                df['Height (m)'][num_PL + num_liquid*i + 6] = step.cyl_f[0] + 2 * step.dome_f[0]                    # Fuel Insulation
                df['Height (m)'][num_PL + num_liquid*i + 7] = " "                                                   # Fuel Residual
                df['Height (m)'][num_PL + num_liquid*i + 8] = step.intertank[0]                                     # Intertank
                df['Height (m)'][num_PL + num_liquid*i + 9] = step.dome_ox[0]                                       # Ox Dome Top
                df['Height (m)'][num_PL + num_liquid*i + 10] = step.cyl_ox[0]                                       # Ox Cylinder
                df['Height (m)'][num_PL + num_liquid*i + 11] = step.dome_ox[0]                                      # Ox Dome Bottom
                df['Height (m)'][num_PL + num_liquid*i + 12] = step.cyl_ox[0] + 2 * step.dome_ox[0]                 # Ox Insulation
                df['Height (m)'][num_PL + num_liquid*i + 13] = " "                                                  # Ox Residual
                df['Height (m)'][num_PL + num_liquid*i + 14] = step.aft_skirt[0]                                    # Aft skirt
                df['Height (m)'][num_PL + num_liquid*i + 15] = step.T_struct                                        # Thrust Struct
                df['Height (m)'][num_PL + num_liquid*i + 16] = " "                                                  # Gimbals
                df['Height (m)'][num_PL + num_liquid*i + 17] = step.L_n                                             # Engine
                df['Height (m)'][num_PL + num_liquid*i + 18] = step.cyl_f[0] + 2 * step.dome_f[0]                   # Fuel (propellant)
                df['Height (m)'][num_PL + num_liquid*i + 19] = step.cyl_ox[0] + 2 * step.dome_ox[0]                 # Ox (propellant)
                
            elif step.propulsion == 'Solid': # MUST STILL ADD SOLID PROPELLANT PARAMETERS IN STEPSIZING FUNCTION (step.srm_casing = [length, SA])
                if i < len(self.listOfSteps) - 1: # if not the last step
                    df['Height (m)'][num_PL + num_solid*i] = step.interstage[0] # Interstage 1 (still called fwd skirt)
                
                elif i == len(self.listOfSteps) - 1: #else if last step
                    df['Height (m)'][num_PL + num_solid*i] = step.fwd_skirt[0] # fwd skirt

                df['Height (m)'][num_PL + num_solid*i + 1] = " "                                                   # Avionics
                df['Height (m)'][num_PL + num_solid*i + 2] = step.total_length                                     # Wiring (Total length of step)
                df['Height (m)'][num_PL + num_solid*i + 3] = step.srm_casing[0]                                    # Solid Propellant Casing
                df['Height (m)'][num_PL + num_solid*i + 4] = " "                                                   # Solid Propellant Residual
                df['Height (m)'][num_PL + num_solid*i + 5] = step.aft_skirt[0]                                     # Aft skirt
                df['Height (m)'][num_PL + num_solid*i + 6] = step.T_struct                                         # Thrust Struct
                df['Height (m)'][num_PL + num_solid*i + 7] = " "                                                   # Gimbals
                df['Height (m)'][num_PL + num_solid*i + 8] = step.L_n                                              # Nozzle
                df['Height (m)'][num_PL + num_solid*i + 9] = step.srm_casing[0]                                    # Solid Propellant

                
        print(df)
       
    def initJ0s():
         pass

    def initThicknesses(self, df, num_PL, num_liquid, num_solid):
        pass
        df['Thickness (m)'][0] = 0.001 # fairing 
        df['Thickness (m)'][1] = " " # payload 
        df['Thickness (m)'][2] = " " # PAF 
        for i in range(len(self.listOfSteps)):
            if self.listOfSteps[i].propulsion == 'Liquid':
                if i < len(self.listOfSteps) - 1:
                    df['Thickness (m)'][num_PL + num_liquid*i] = 0.001 # Interstage
                
                elif i == len(self.listOfSteps) - 1:
                    df['Thickness (m)'][num_PL + num_liquid*i] = 0.001                  # Forward Skirt

                df['Thickness (m)'][num_PL + num_liquid*i + 1] = " "                    # Avionics
                df['Thickness (m)'][num_PL + num_liquid*i + 2] = " "                    # Wiring
                df['Thickness (m)'][num_PL + num_liquid*i + 3] = 0.001                  # Fuel Dome Top
                df['Thickness (m)'][num_PL + num_liquid*i + 4] = 0.001                  # Fuel Cylinder
                df['Thickness (m)'][num_PL + num_liquid*i + 5] = 0.001                  # Fuel Dome Bottom
                df['Thickness (m)'][num_PL + num_liquid*i + 6] = " "                    # Fuel Insulation
                df['Thickness (m)'][num_PL + num_liquid*i + 7] = " "                    # Fuel Residual
                df['Thickness (m)'][num_PL + num_liquid*i + 8] = 0.001                  # Intertank
                df['Thickness (m)'][num_PL + num_liquid*i + 9] = 0.001                  # Ox Dome Top
                df['Thickness (m)'][num_PL + num_liquid*i + 10] = 0.001                 # Ox Cylinder
                df['Thickness (m)'][num_PL + num_liquid*i + 11] = 0.001                 # Ox Dome Bottom
                df['Thickness (m)'][num_PL + num_liquid*i + 12] = " "                   # Ox Insulation
                df['Thickness (m)'][num_PL + num_liquid*i + 13] = " "                   # Ox Residual
                df['Thickness (m)'][num_PL + num_liquid*i + 14] = 0.001                 # Aft skirt
                df['Thickness (m)'][num_PL + num_liquid*i + 15] = " "                   # Thrust Struct
                df['Thickness (m)'][num_PL + num_liquid*i + 16] = " "                   # Gimbals
                df['Thickness (m)'][num_PL + num_liquid*i + 17] = " "                   # Engine
                df['Thickness (m)'][num_PL + num_liquid*i + 18] = " "                   # Fuel (propellant)
                df['Thickness (m)'][num_PL + num_liquid*i + 19] = " "                   # Ox (propellants
            elif self.listOfSteps[i].propulsion == 'Solid':
                df['Thickness (m)'][num_PL + num_solid*i + 1] = " "                     # Avionics
                df['Thickness (m)'][num_PL + num_solid*i + 2] = " "                     # Wiring
                df['Thickness (m)'][num_PL + num_solid*i + 3] = " "                     # Solid Propellant Casing
                df['Thickness (m)'][num_PL + num_solid*i + 4] = " "                     # Solid Propellant Residual
                df['Thickness (m)'][num_PL + num_solid*i + 5] = 0.001                   # Aft skirt
                df['Thickness (m)'][num_PL + num_solid*i + 6] = " "                     # Thrust Struct
                df['Thickness (m)'][num_PL + num_solid*i + 7] = " "                     # Gimbal
                df['Thickness (m)'][num_PL + num_solid*i + 8] = " "                     # Nozzle
                df['Thickness (m)'][num_PL + num_solid*i + 9] = " "                     # Solid Propellant

    def initMasses(self, df, num_PL, num_liquid, num_solid):
        df['Mass (kg)'][0] = self.listOfSteps[len(self.listOfSteps)-1].fairing[1] * self.listOfSteps[len(self.listOfSteps)-1].rho_fairing * df['Thickness (m)'][0] # Fairing
        df['Mass (kg)'][1] = self.PL                                                                                                                               # Payload
        df['Mass (kg)'][2] = self.listOfSteps[len(self.listOfSteps)-1].m_PL - self.PL                                                                              # PAF

        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            if step.propulsion == 'Liquid':
                if i < len(self.listOfSteps) - 1: # if not last step
                    df['Mass (kg)'][num_PL + num_liquid*i] = step.interstage[1] * self.rho_body * df['Thickness (m)'][num_PL + num_liquid*i]                        # Interstage 1
                
                elif i == len(self.listOfSteps) - 1: # if last step
                    df['Mass (kg)'][num_PL + num_liquid*i] = step.fwd_skirt[1] * self.rho_body * df['Thickness (m)'][num_PL + num_liquid*i]                         # Forward Skirt
                    #print("fwd skirt has mass " + str(step.fwd_skirt[1] * self.rho_body * df['Thickness (m)'][num_PL + num_liquid*i]) + " kg")

                df['Mass (kg)'][num_PL + num_liquid*i + 1] =  10*pow(self.m_0[i], 0.361) # Avionics
                #print("Avionics has mass " + str(10*pow(self.m_0[i], 0.361)) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 2] = 1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4) # Wiring 
                #print("Wiring has mass " + str(1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4)) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 3] = step.dome_f[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 3]                          # Fuel Dome Top
                df['Mass (kg)'][num_PL + num_liquid*i + 4] = step.cyl_f[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 4]                           # Fuel Cylinder
                df['Mass (kg)'][num_PL + num_liquid*i + 5] = step.dome_f[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 5] # Fuel Dome Bottom
                df['Mass (kg)'][num_PL + num_liquid*i + 6] = (step.cyl_f[1] + 2 * step.dome_f[1]) * step.SA_rho_insulation['Fuel']                                    # Fuel Insulation
                #print("Fuel Insulation has mass " + str((step.cyl_f[1] + 2 * step.dome_f[1]) * step.SA_rho_insulation['Fuel']) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 7] = step.residual_prop_perc * step.m_f_ideal                                                                 # Fuel Residual
                #print("Fuel Residual has mass " + str(step.residual_prop_perc * step.m_f_ideal) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 8] = step.intertank[1] * self.rho_body * df['Thickness (m)'][num_PL + num_liquid*i + 8]                       # Intertank
                print("Intertank has mass " + str(step.intertank[1] * self.rho_body * df['Thickness (m)'][num_PL + num_liquid*i + 8]) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 9] = step.dome_ox[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 9]                         # Ox Dome Top
                df['Mass (kg)'][num_PL + num_liquid*i + 10] = step.cyl_ox[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 10]                         # Ox Cylinder
                df['Mass (kg)'][num_PL + num_liquid*i + 11] = step.dome_ox[1] * step.rho_tank * df['Thickness (m)'][num_PL + num_liquid*i + 11]                        # Ox Dome Bottom
                df['Mass (kg)'][num_PL + num_liquid*i + 12] = (step.cyl_ox[1] + 2 * step.dome_ox[1]) * step.SA_rho_insulation['Oxidizer']                             # Ox Insulation
                #print("Ox Insulation has mass " + str((step.cyl_ox[1] + 2 * step.dome_ox[1]) * step.SA_rho_insulation['Oxidizer']) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 13] = step.residual_prop_perc * step.m_ox_ideal                                                               # Ox Residual
                #print("Ox Residual has mass " + str(step.residual_prop_perc * step.m_ox_ideal) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 14] = step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][num_PL + num_liquid*i + 14]                    # Aft skirt
                #print("Aft Skirt has mass " + str(step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][num_PL + num_liquid*i + 10]) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 15] = 2.55*pow(10, -4)*step.T_SL                                                                              # Thrust Struct
                #print("Thrust Structure has mass " + str(2.55*pow(10, -4)*step.T_SL) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 16] = step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)                             # Gimbals
                #print("Nmber of gimballed engines is " + str(step.num_gimballed_engines))
                print("Total gimbals mass is " + str(step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 17] = step.num_engines * (step.T_SL_engine * (7.81 * pow(10, -4) + 3.37 * pow(10, -5) * sqrt(step.epsilon)) + 59) # Engines
                print("Total engines mass is " + str(step.T_SL * (7.81 * pow(10, -4) + 3.37 * pow(10, -5) * sqrt(step.epsilon)) + 59) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 18] = step.m_f_ideal + step.fuel_frac * step.startup_prop                                                     # Fuel (propellant)
                #print("Total Fuel mass is " + str(step.m_f_ideal + step.fuel_frac * step.startup_prop ) + " kg")
                df['Mass (kg)'][num_PL + num_liquid*i + 19] = step.m_ox_ideal + step.ox_frac * step.startup_prop                                                      # Ox (propellant)
                print("Total Ox mass is " + str(step.m_ox_ideal + step.ox_frac * step.startup_prop) + " kg") # take in thicknesses as a parameter in order to iterate various thicknesses

            elif step.propulsion == 'Solid':
                df['Mass (kg)'][num_PL + num_solid*i + 1] =  10*pow(self.m_0[i], 0.361) # Avionics
                #print("Avionics has mass " + str(10*pow(self.m_0[i], 0.361)) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 2] = 1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4) # Wiring 
                #print("Wiring has mass " + str(1.058*pow(self.m_0[i], 1/2) * pow(step.total_length, 1/4)) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 3] = step.residual_prop_perc * step.m_f_ideal                                                                 # Fuel Residual
                #print("Fuel Residual has mass " + str(step.residual_prop_perc * step.m_f_ideal) + " kg")
                #print("Ox Insulation has mass " + str((step.cyl_ox[1] + 2 * step.dome_ox[1]) * step.SA_rho_insulation['Oxidizer']) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 4] = step.residual_prop_perc * step.m_ox_ideal                                                               # Ox Residual
                #print("Ox Residual has mass " + str(step.residual_prop_perc * step.m_ox_ideal) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 5] = step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][num_PL + num_solid*i + 5]                    # Aft skirt
                #print("Aft Skirt has mass " + str(step.aft_skirt[1] * self.rho_body *  df['Thickness (m)'][num_PL + num_solid*i + 10]) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 6] = 2.55*pow(10, -4)*step.T_SL                                                                              # Thrust Struct
                #print("Thrust Structure has mass " + str(2.55*pow(10, -4)*step.T_SL) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 7] = step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)                             # Gimbals
                #print("Nmber of gimballed engines is " + str(step.num_gimballed_engines))
                print("Total gimbals mass is " + str(step.num_gimballed_engines * 237.8*pow(step.T_SL_engine / step.p_c, 0.9375)) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 8] = step.num_engines * (step.T_SL_engine * (7.81 * pow(10, -4) + 3.37 * pow(10, -5) * sqrt(step.epsilon)) + 59)  # Engines
                #print("Total engines mass is " + str(step.T_SL * (7.81 * pow(10, -4) + 3.37 * pow(10, -5) * sqrt(step.epsilon)) + 59) + " kg")
                df['Mass (kg)'][num_PL + num_solid*i + 9] = self.m_p[i] + step.startup_prop                                                                         # Solid Propellant
                #print("Total Fuel mass is " + str(step.m_f_ideal + step.fuel_frac * step.startup_prop ) + " kg")
    def initDistances(self, df, num_PL, num_liquid, num_solid): # initiate distances from bottom of aft skirt to CM of component
        pass
        # Note the distances are initialized in reverse order, but maintain their order in the dataframe df
        pi = 3.1415926535897932
        for i in range(len(self.listOfSteps)):
            step = self.listOfSteps[i]
            df['Distance (m)'][num_PL + num_liquid*i + 19] = df['Distance (m)'][num_PL + num_liquid*i + 10] + step.cyl_ox[0]/2                                          # Ox (propellant)
            df['Distance (m)'][num_PL + num_liquid*i + 18] = df['Distance (m)'][num_PL + num_liquid*i + 4]                                                              # Fuel (propellant)
            
            if i == 0: # if first step
                df['Distance (m)'][num_PL + num_liquid*i + 17] = step.r/4                                                                                                   # Engine
                df['Distance (m)'][num_PL + num_liquid*i + 16] = step.L_n                                                                                                   # Gimbals
                df['Distance (m)'][num_PL + num_liquid*i + 15] = step.T_struct/2                                                                                            # Thrust Struct
                df['Distance (m)'][num_PL + num_liquid*i + 14] = step.aft_skirt[0]/2                                                                                        # Aft skirt
            else:
                df['Distance (m)'][num_PL + num_liquid*i + 17] = df['Distance (m)'][num_PL + num_liquid*(i-1) + 4] + self.listOfSteps[i-1].cyl_f[0] + self.listOfSteps[i-1].interstage[0] + step.r/4                                                                                                   # Engine
                df['Distance (m)'][num_PL + num_liquid*i + 16] = df['Distance (m)'][num_PL + num_liquid*(i-1) + 4] + self.listOfSteps[i-1].cyl_f[0] + self.listOfSteps[i-1].interstage[0] + step.L_n                                                                                                   # Gimbals
                df['Distance (m)'][num_PL + num_liquid*i + 15] = df['Distance (m)'][num_PL + num_liquid*(i-1) + 4] + self.listOfSteps[i-1].cyl_f[0] + self.listOfSteps[i-1].interstage[0] + step.T_struct/2                                                                                            # Thrust Struct
                df['Distance (m)'][num_PL + num_liquid*i + 14] = df['Distance (m)'][num_PL + num_liquid*(i-1) + 4] + self.listOfSteps[i-1].cyl_f[0] + self.listOfSteps[i-1].interstage[0] + step.aft_skirt[0]/2 #  Aft skirt = previous fwd skirt top + 1/2 * current step aft skirt height
            df['Distance (m)'][num_PL + num_liquid*i + 13] = " "                                                                                                        # Ox Residual
            df['Distance (m)'][num_PL + num_liquid*i + 12] = df['Distance (m)'][num_PL + num_liquid*i + 14] + step.aft_skirt[0]/2 + step.cyl_ox[0]/2                    # Ox Insulation
            df['Distance (m)'][num_PL + num_liquid*i + 11] = df['Distance (m)'][num_PL + num_liquid*i + 14] + step.r/2 + (step.dome_ox[0]- 4 * step.dome_ox[0]/(3 * pi))# Ox Dome Bottom
            df['Distance (m)'][num_PL + num_liquid*i + 10] = df['Distance (m)'][num_PL + num_liquid*i + 14] + step.aft_skirt[0]/2 + step.cyl_ox[0]/2                    # Ox Cylinder
            df['Distance (m)'][num_PL + num_liquid*i + 9] = df['Distance (m)'][num_PL + num_liquid*i + 10] + step.cyl_ox[0]/2 + 4 * step.dome_ox[0]/(3 * pi)            # Ox Dome Top
            df['Distance (m)'][num_PL + num_liquid*i + 8] = df['Distance (m)'][num_PL + num_liquid*i + 10] + step.cyl_ox[0]/2 + step.dome_ox[0] + step.r/4              # Intertank
            df['Distance (m)'][num_PL + num_liquid*i + 7] = " "                                                                                                         # Fuel Residual
            df['Distance (m)'][num_PL + num_liquid*i + 6] = df['Distance (m)'][num_PL + num_liquid*i + 8] + step.intertank[0]/2 + step.cyl_f[0]/2                       # Fuel Insulation
            df['Distance (m)'][num_PL + num_liquid*i + 5] = df['Distance (m)'][num_PL + num_liquid*i + 8] + step.intertank[0]/2 - 4 * step.dome_f[0]/(3 * pi)           # Fuel Dome Bottom
            df['Distance (m)'][num_PL + num_liquid*i + 4] = df['Distance (m)'][num_PL + num_liquid*i + 8] + step.intertank[0]/2 + step.cyl_f[0]/2                       # Fuel Cylinder
            df['Distance (m)'][num_PL + num_liquid*i + 3] = df['Distance (m)'][num_PL + num_liquid*i + 4] + step.cyl_f[0]/2 + 4 * step.dome_f[0]/(3 * pi)               # Fuel Dome Top
            df['Distance (m)'][num_PL + num_liquid*i + 2] = step.total_length/2                                                                                         # Wiring
            df['Distance (m)'][num_PL + num_liquid*i + 1] = " "                                                                                                         # Avionics

            
            if i < len(self.listOfSteps) - 1: # if not last step
                #df['Distance (m)'][num_PL + num_liquid*i] = df['Distance (m)'][num_PL + num_liquid*(i-1)] # Interstage 1
                 df['Distance (m)'][num_PL + num_liquid*i] = df['Distance (m)'][num_PL + num_liquid*i + 4] + step.cyl_f[0]/2 + step.interstage[0]/3*(2*self.listOfSteps[i+1].r + step.r)/(self.listOfSteps[i+1].r + step.r) # Interstage ( using formual for centroid of trapezoid)
            elif i == len(self.listOfSteps) - 1: # if last step
                 df['Distance (m)'][num_PL + num_liquid*i] = df['Distance (m)'][num_PL + num_liquid*i + 3] - step.cyl_f[0]/2 + step.fwd_skirt[0]/2 # fwd skirt
        
        df['Distance (m)'][2] = df['Distance (m)'][num_PL + num_liquid*(self.num_steps-1)] # PAF
        df['Distance (m)'][1] = df['Distance (m)'][2] + df['Height (m)'][1]/2 # payload
        df['Distance (m)'][0] = df['Distance (m)'][num_PL + num_liquid*(self.num_steps-1)] + 1/3*df['Height (m)'][0] # fairing

    def AddSlide(self):
        prs = Presentation() # Create Presentation

        Vehicle_slide_layout = prs.slide_layouts[5] # define "title only" layout
        vehicle_diagram = prs.slides.add_slide(Vehicle_slide_layout) # add vehicle diagram slide
        shapes = vehicle_diagram.shapes # declare shapes for ease of access to the slide's shapes
        shapes.title.text = self.name + ' Diagram'
        
        scale = 0.1
        #for i in self.listOfSteps:
        #    aft_skirt = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(, top1, width1, height1)

        prs.save(self.name + 'Diagram.pptx')



class Step(LaunchVehicle):
    pass

    def __init__(self, LV, r, step_num, TW, engine, dome_shape, propellants, insulation_dict, tank_material, num_engines, num_gimballed_engines, t_start, fairing_material, fairing_shape, ):
        pass
        self.LV = LV
        if step_num == LV.num_steps:
            self.m_PL = 1.15*LV.PL + 15
            self.r = r
            self.step_num = step_num
            self.TW = TW
            self.engine = engine
            self.dome_shape = dome_shape
            self.propellants = propellants
            self.insulation_dict = insulation_dict
            self.tank_material = tank_material
            self.num_engines = num_engines
            self.num_gimballed_engines = num_gimballed_engines
            self.fairing_material = fairing_material
            self.fairing_shape = fairing_shape
            self.t_start = t_start
        else:
            self.m_PL = 0
            self.r = r
            self.step_num = step_num
            self.TW = TW
            self.engine = engine
            self.dome_shape = dome_shape
            self.propellants = propellants
            self.insulation_dict = insulation_dict
            self.tank_material = tank_material
            self.num_engines = num_engines
            self.num_gimballed_engines = num_gimballed_engines
            self.fairing_material = 0
            self.fairing_shape = 0
            self.t_start = t_start

    def print(self):
        print("This is step " + str(self.step_num) + " of " + self.LV.name)
        print("The radius is " + str(self.r))
        print("The Thrust-to-Weight is "+ str(self.TW))
        print("The engine is based on the " + self.engine)
        print("The dome shape is " + self.dome_shape)
        print("The propellant combination is " + self.propellants)
        print("The body material is " + self.LV.body_material)
        print("The tank material is " + self.tank_material)
        print("The number of engines is " + str(self.num_engines))
        if self.step_num == self.LV.num_steps:
            print("The fairing material is " + self.fairing_material)
            print("The fairing shape is " + "'" + self.fairing_shape + "'")

    def sizeStep(self):
    ## Assumptions
    # Elliptical-2 dome has AR = 2 and Elliptical-sqr2 has AR = sqrt(2)
    # thickness of fairing is 0.005 m
    # Bulkheads are uncommon
    # volume of single cubesat is 0.001 m^3
    # mass of single cubeSat is 1.33 kg
    # payload volume margin is 1.5: accounts for air gaps between cubeSats and
    # fairing wall
    #STEPSIZING Summary of this function goes here
    # The inputs of this function are:
        # if this step is the last/upper => last_step = true, otherwise it is false
        # the gross lift-mass in kg => m_gross in kg, e.g. 6414.5
        # TW - Thrust to weight ratio for the step
        # m_prop =  propellant mass(kg)
        # the engine used to size engine_params => engine = "Raptor" or
        # "Merlin"
        # type of propellant tank dome => dome = "Elliptical-2" (hemispheres) or "Elliptical-sqr2" or
        # "Circular" (hemispheres) or "Spherical"
        # type of propellants => propellants = "Methalox" or "Kerolox"
        # material of the tanks for shrinkage => right now the only option
        # is => tank_material = "Aluminum 6061"
        # number of engines for this step => num_engines = 1, 2, 3, 4, etc...
        # the material of the fairing is fairing_material = "Aluminum 6061",
        # for example
        # the shape of the fairing is fairing_shape = "Cone", "Ogive", etc.
        # the payload mass is m_PL, INCLUDING the PL attach fitting
        #self.step_params = {'Radius': None, 'Exit diameter': None, 'Nozzle length': None, 'Thrust structure': None, 'Aft skirt': None, 'Dome_f': None, 'Cyl_f': None, 'Intertank': None, 'Dome_ox': None, 'Cyl_ox': None, 'Fairing': None, 'Fwd skirt': None}
        self.aft_skirt = []
        self.dome_f = []
        self.cyl_f = []
        self.intertank = []
        self.dome_ox = []
        self.cyl_ox = []
        self.srm_casing = []
        if self.step_num == self.LV.num_steps:
            self.fairing = []
            self.fwd_skirt = []

        # SET PARAMETERS
        pi = 3.1415926535897932#3846264338327950288
        g = 9.80665  # gravitational acceleration (m/s^2)
        if self.engine == "Raptor":
            Isp = 330  # specific impulse (s) ACTUAL
            #Isp = 365 # specific impulse (s) ON EXCEL 
            Isp_vac = 380  # vacuum specific impulse (s)
            self.p_c = 30*pow(10,6)  # chamber pressure (Pa) ACTUAL
            #p_c = 9.7*pow(10,6)  # chamber pressure (Pa) ON EXCEL
            self.epsilon = 45  # expansion ratio

        elif self.engine == "Merlin":
            Isp = 282  # specific impulse (s)
            Isp_vac = 311  # vacuum specific impulse (s)
            self.p_c = 9.7*pow(10,6)  # chamber pressure (Pa)
            self.epsilon = 22  # expansion ratio
     
        if self.dome_shape == "Elliptical-2":
            AR = 2  # aspect ratio ( > 1 means elliptical hemisphere, = 1 means spherical hemisphere) (1)
        elif self.dome_shape == "Elliptical-sqr2":
            AR = sqrt(2) 
        elif self.dome_shape == "Circular":
            AR = 1 

        self.SA_rho_insulation = {'Oxidizer': 0, 'Fuel': 0}
        if self.insulation_dict['Fuel'] == 'Rubber':
            self.SA_rho_insulation['Fuel'] = 0.02
        if self.insulation_dict['Oxidizer'] == 'Rubber':
            self.SA_rho_insulation['Oxidizer'] = 0.02

        
        if self.propellants == "Kerolox":
            self.propulsion = 'Liquid'
            self.rho_f = 800  # fuel densitity (kg/m^3)
            self.rho_ox = 1140   # oxidizer density (kg/m^3)
            OF = 2.34  # O/F mixture ratio (1)
            FO = 1/OF  # F/O mixture ratio (1)
            self.fuel_frac = FO/(1 + FO)  # fraction of fuel mass (1)
            
        elif self.propellants == "Methalox":
            self.propulsion = 'Liquid'
            self.rho_f = 424  # fuel densitity (kg/m^3)
            self.rho_ox = 1140   # oxidizer density (kg/m^3)
            OF = 3.55  # O/F mixture ratio (1)
            FO = 1/OF  # F/O mixture ratio (1)
            self.fuel_frac = FO/(1 + FO)  # fraction of fuel mass (1)

        elif self.propellants == "AP-Al-HTPB":
            self.propulsion = 'Solid'
            self.rho_prop = 1600 # propellant density (kg/m^3) https://www.sciencedirect.com/science/article/pii/S1878535215000106#bb0020

        rm_temp = 293  # room temperature (K)
        if self.tank_material == "Aluminum 6061":
            a = 23.4 * pow(10, -6)  # (1/K)
            self.rho_tank = 2700  # kg/m^3
        if self.fairing_material == "Aluminum 6061":
                self.rho_fairing = 2700  # kg/m^3 aluminum density

        if self.propellants == "Methalox":
            fuel_temp = 100  # current temperature of fuel(K)
            ox_temp =80  # boiling temp of ox (K)
        elif self.propellants == "Kerolox":
            fuel_temp = 293 
            ox_temp = 80 

        #step = 2 
        # Parameters for actual tank volumes
        step_dia = 2*self.r  # diameter (m)
        self.circumference = 2*pi*self.r  # m Circumference
      
        # Propellant Mass Calculations
    
        self.residual_prop_perc = 0.02  # residual propellant (1)
        self.T_SL = 0
        for i in range(len(self.LV.m_0))[::-1]:
            self.T_SL += self.LV.m_0[i]*g
            if self.step_num-1 == i:
                break
            #if step_num - 1 == 
        #self.T_SL = self.TW * self.LV.m_0[self.step_num-1] 
        self.T_SL_engine = self.T_SL/self.num_engines
        #disp(T_SL)
        residual_prop = self.residual_prop_perc * self.LV.m_p[self.step_num-1]  # residual propellant (kg)
        m_dot = self.T_SL / (Isp *g)  # mass flow rate (kg/s)
        self.startup_prop = m_dot*self.t_start  # propellant needed for start-up (kg)
        m_prop_tot = self.LV.m_p[self.step_num-1] + residual_prop + self.startup_prop 

        ## Needed Tank Volumes
        if self.propulsion == 'Liquid':
            self.m_f_ideal = m_prop_tot* self.fuel_frac
            self.vol_f_ideal = self.m_f_ideal/self.rho_f  # ideal fuel volume (kg)
            self.ox_frac = 1/(1+FO) 
            self.m_ox_ideal = m_prop_tot * self.ox_frac
            self.vol_ox_ideal = self.m_ox_ideal/self.rho_ox  # ideal ox volume (kg)
            # #disp(self.vol_f_ideal)
            # #disp(self.vol_ox_ideal)

            ## Shrinkage
            # Shrinkage Parameters
            shrinkage_fuel = 3*a*(rm_temp - fuel_temp)  # shrinkage factor for fuel tank (1)
            shrinkage_ox = 3*a*(rm_temp - ox_temp)  # shrinkage factor for ox tank (1)
            #disp(shrinkage_ox)
            ## Total Tank Volumes
            # Tank Sizing Parameters
            self.ullage_frac = 0.025  # ullage factor (how much gas is in the propellant tank) (1)
            vol_tot_needed_f = self.vol_f_ideal * (1 + self.ullage_frac + shrinkage_fuel)
            vol_tot_needed_ox = self.vol_ox_ideal * (1 + self.ullage_frac + shrinkage_ox) 
            self.m_tot_needed_f = vol_tot_needed_f * self.rho_f
            self.m_tot_needed_ox = vol_tot_needed_ox * self.rho_ox
            #disp(vol_tot_needed_f)
            #disp(vol_tot_needed_ox)

            ## Actual Tank Volumes Calculations
            e = sqrt(1 - pow(1/AR,2))  # eccentricity (1)

            ## Fuel
            self.dome_f.append(self.r/AR)  # m Lengths of fuel dome (m)
            #disp(self.dome_f[0])
            self.dome_f.append(pi*pow(self.r,2)*( 1 + 1/(2*e*pow(AR,2))*log((1+e)/(1-e))))  # m^2 SA of the fuel dome (m^2)
            dome_vol_f = 2/3*pi*pow(self.r,2)*self.dome_f[0]  # m^3 dome fuel tank volume (m^3)
            self.dome_f.append(dome_vol_f)
            #disp(dome_vol_f)
            cyl_vol_f = vol_tot_needed_f - 2*dome_vol_f  # m^3 cylinder fuel tank volume (m^3)
            self.cyl_f.append(cyl_vol_f/(pi*pow(self.r,2)))  # m Lengths of the fuel cylinder (m)
        
            # Iterate sizeStep, changing step diameter to avoid a negative length of tank cylinder
            if self.cyl_f[0] < 0:
                self.r=self.r-0.01
                self.sizeStep() 
                return
     
            self.cyl_f.append(self.cyl_f[0]*self.circumference)  # m^2 SA of fuel cylinder
            #disp(self.self.cyl_f[0])
            ## Oxidizer
            h_dome_ox = self.r/AR # Lengths of the oxidizer dome (m)
            self.dome_ox.append(h_dome_ox)  
            #disp(self.dome_ox[0])
            SA_dome_ox = pi*pow(self.r,2)*( 1 + 1/(2*e*pow(AR,2))*log((1+e)/(1-e))) # SA of the oxidizer dome (m^2)
            self.dome_ox.append(SA_dome_ox)  
            dome_vol_ox = 2/3*pi*pow(self.r,2)*self.dome_ox[0]  # dome oxidizer tank volume (m^3)
            self.dome_ox.append(dome_vol_ox)
            cyl_vol_ox = vol_tot_needed_ox - 2*dome_vol_ox  # cylinder fuel tank volume (m^3)
            h_cyl_ox = cyl_vol_ox/(pi*pow(self.r,2)) # length of oxidizer cylinder
            self.cyl_ox.append(h_cyl_ox)  # append to list the length of the ox cylinder
            #disp(self.cyl_ox[0])

            # Iterate sizeStep, changing step diameter to avoid a negative length of tank cylinder
            if self.cyl_ox[0] < 0:
                self.r=self.r-0.01
                self.sizeStep()
                return
            SA_cyl_ox = h_cyl_ox*self.circumference
            self.cyl_ox.append(SA_cyl_ox)  # m^2 SA of oxidizer cylinder
            self.cyl_ox.append(cyl_vol_ox) # m^3 of oxidizer cylinder

            self.intertank.append(1/4*step_dia+self.dome_f[0]+self.dome_ox[0])  # m length of intertank
            self.intertank.append(self.intertank[0]*self.circumference)   # m^2 surface area of intertank
        elif self.propulsion == 'Solid':
            # Ideal volumes
            self.m_prop_ideal = m_prop_tot
            self.vol_prop_ideal = self.m_prop_ideal / self.rho_prop

            # NO Shrinkage
            # NO ULLAGE
            burn_pattern = 0.15 # assume 15% of casing space is allocated for burn pattern
            vol_tot_needed_prop = self.vol_prop_ideal * (1 + burn_pattern)
            self.m_tot_needed_prop = vol_tot_needed_prop * self.rho_prop
            h_srm_casing = vol_tot_needed_prop / (pi * pow(self.r,2))
            self.srm_casing.append(h_srm_casing)
            SA_srm_casing = h_srm_casing * pi * step_dia
            self.srm_casing.append(SA_srm_casing)

            
        ## Liquid Bipropellant Engine Sizing
        Lth_ratio = 0.9  # chamber length / chamber diameter ratio
        alpha_conv = radians(30)  # convergent section angle (radians)
        L_star = 0.75  # characteristic length (m)
        p_inf = 101.3*pow(10,3)  # ambient pressure (Pa)
        if self.step_num == self.LV.num_steps:
            T_vac = self.T_SL 
        else:
            T_vac = self.T_SL/(1-(p_inf*Isp*g*self.epsilon)/(self.p_c*g*Isp_vac))  # total vacuum thrust ACTUAL
            #T_vac = self.T_SL/(1-(p_inf*Isp*g*self.epsilon)/(self.p_c*g*311))  # total vacuum thrust ON EXCEL


        T_vac_engine = T_vac/self.num_engines  # vacuum thrust of each engine
        A_t = Isp*g*T_vac_engine/(self.p_c*g*Isp_vac)  # throat area ACTUAL
        #A_t = Isp*g*T_vac_engine/(self.p_c*g*Isp_vac*2)  # throat area ON EXCEL

        d_t = sqrt(4*A_t/pi)  # throat diameter
        self.d_e = self.epsilon*d_t  # exit diameter
        self.L_n = (self.d_e-d_t)*0.8/((2*tan(radians(30))))  # length of nozzle for 80# length & 15 deg cone
        #print(self.L_n)
        L_c = pow((4*L_star * A_t * Lth_ratio)/pi, 1/3)  # chamber length
        d_c=  L_c/Lth_ratio 
        L_conv = (d_c - d_t)/(2*tan(alpha_conv)) 
        #self.L_eng_tot = L_n + L_c + L_conv 
        #disp(L_conv)
        ## Intertank and Thrust Structure Sizing
        # The following variables denote length
        
      
        if self.step_num == self.LV.num_steps:
            self.fwd_skirt.append(1/3*step_dia + self.dome_ox[0])  # m Length of fwd skirt
            self.fwd_skirt.append(self.fwd_skirt[0]*self.circumference)  # m^2 SA of fwd skirt
            self.aft_skirt.append(self.fwd_skirt[0])  # m Length of aft skirt
            self.aft_skirt.append(self.aft_skirt[0]*self.circumference)  # m^2 SA of aft skirt
            # NOTE: Change T_struct based on nozzle length
            self.T_struct = self.aft_skirt[0]/2  # length of thrust structure
            
         
            # NEEDED Payload VOLUME
            cubeSat_vol = 0.001  # m^3 volume of individual cubeSat
            m_cubeSat = 1.33  # kg mass of individual cubeSat
            rho_cubeSat = m_cubeSat/cubeSat_vol  # kg/m^3 density of cubeSat 
            PL_vol_needed = 1.5 * self.m_PL/rho_cubeSat  # m^3 needed payload volume (FS = 1.5)
            #disp(PL_vol_needed)
            # ACTUAL PAYLOAD VOLUME
            if self.fairing_shape == "Cone":
                self.fairing.append(3* PL_vol_needed / (pi *pow(self.r,2))) 
         
            #disp(self.L_fairing)
                self.fairing.append(pi * self.r * ( self.r + sqrt( pow(self.fairing[0],2) + pow(self.r,2) )))  # m^2 fairing surface area
            #disp(m_fairing)
            #disp(SA_fairing)
        
        else:
            self.aft_skirt.append(step_dia/2 + 3/4*self.L_n)  # m length of the aft skirt
            self.aft_skirt.append(self.aft_skirt[0]*self.circumference) 
            # NOTE: Change T_struct based on nozzle length
            self.T_struct = self.aft_skirt[0]/2  # m length of thrust structure
      
        
        #print(self.intertank)
        #disp(self.intertank[0])
        # get total length of step
        if self.step_num == self.LV.num_steps:
            if self.propulsion == 'Liquid':
                self.total_length = self.aft_skirt[0] + self.fwd_skirt[0] + self.intertank[0] + self.cyl_f[0] + self. cyl_ox[0] 
            elif self.propulsion == 'Solid':
                self.total_length = self.aft_skirt[0] + self.fwd_skirt[0] + self.srm_casing[0]
        else:
            if self.propulsion == 'Liquid':
                self.total_length = self.aft_skirt[0] + self.intertank[0] + self.cyl_f[0] + self. cyl_ox[0] 
            elif self.propulsion == 'Solid':
                self.total_length = self.aft_skirt[0] + self.srm_casing[0]
        ## Setting fcn Output Variables
        #lengths_key = ["self.d_e", "self.L_n", "r", "self.L_T_struct", "self.aft_skirt[0]", "self.cyl_ox[0]", "self.dome_ox[0]", "self.intertank[0]", "self.cyl_f[0]", "self.dome_f[0]", "self.fwd_skirt[0]", "self.fairing[0]"] 
        #if self.step_num == self.LV.num_steps:
        #    pass
        #    #self.step_parameters = [self.d_e, self.L_n, self.r, self.L_T_struct, self.aft_skirt[0], self.cyl_ox[0], self.dome_ox[0], self.intertank[0], self.cyl_f[0], self.dome_f[0], self.fwd_skirt[0], self.fairing[0]] 
        #    #self.t_rho = [self.aft_skirt[1]*rho_body, self.dome_ox[1]*self.rho_tank, self.cyl_ox[1]*self.rho_tank, self.intertank[1]*rho_body, self.dome_f[1]*self.rho_tank ,self.cyl_f[1]*self.rho_tank, self.fwd_skirt[1]*rho_body] 
        #else:
        #    pass
        #    #self.step_parameters = [self.d_e, self.L_n, self.r, self.L_T_struct, self.aft_skirt[0], self.cyl_ox[0], self.dome_ox[0], self.intertank[0], self.cyl_f[0], self.dome_f[0]] 
        #    #self.t_rho = [self.aft_skirt[1]*rho_body, self.dome_ox[1]*self.rho_tank, self.cyl_ox[1]*self.rho_tank, self.intertank[1]*rho_body, self.dome_f[1]*self.rho_tank ,self.cyl_f[1]*self.rho_tank, rho_body*self.circumference] 
     
    
    
 
